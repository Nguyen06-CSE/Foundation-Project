# app/services/file_processor.py
from __future__ import annotations

import logging
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

STORAGE_DIR = Path("storage")
THUMBNAIL_DIR = STORAGE_DIR / "thumbnails"
THUMBNAIL_DIR.mkdir(parents=True, exist_ok=True)

# Map MIME type → extension nhóm
MIME_TO_GROUP = {
    "application/pdf": "pdf",
    "application/msword": "docx",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.ms-powerpoint": "pptx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "image/jpeg": "image",
    "image/png": "image",
    "image/webp": "image",
    "image/gif": "image",
    "text/plain": "text",
}


def get_file_group(mime_type: str) -> str:
    return MIME_TO_GROUP.get(mime_type, "other")


def extract_text(file_path: str, mime_type: str) -> str:
    """
    Extract text từ file. Trả về string rỗng nếu lỗi.
    - PDF có text: dùng pdfplumber
    - PDF scan (text rỗng): fallback OCR pytesseract
    - DOCX: python-docx
    - PPTX: python-pptx  
    - Image: OCR pytesseract
    - Text: đọc trực tiếp
    """
    path = Path(file_path)
    if not path.exists():
        logger.warning(f"File không tồn tại: {file_path}")
        return ""

    group = get_file_group(mime_type)

    try:
        if group == "pdf":
            import pdfplumber
            with pdfplumber.open(path) as pdf:
                text = "\n".join(
                    page.extract_text() or "" for page in pdf.pages
                )
            # Nếu PDF scan (không có text layer) → OCR
            if not text.strip():
                logger.info(f"PDF không có text layer, chuyển sang OCR: {file_path}")
                text = _ocr_pdf(path)
            return text[:100_000]  # giới hạn 100k ký tự

        elif group == "docx":
            from docx import Document as DocxDocument
            doc = DocxDocument(path)
            return "\n".join(p.text for p in doc.paragraphs)[:100_000]

        elif group == "pptx":
            from pptx import Presentation
            prs = Presentation(path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts)[:100_000]

        elif group == "image":
            return _ocr_image(path)

        elif group == "text":
            return path.read_text(encoding="utf-8", errors="ignore")[:100_000]

    except Exception as e:
        logger.error(f"extract_text lỗi [{file_path}]: {e}")

    return ""


def create_thumbnail(file_path: str, mime_type: str, doc_id: int) -> Optional[str]:
    """
    Tạo thumbnail JPG. Trả về đường dẫn tương đối hoặc None nếu không hỗ trợ.
    - PDF: render trang đầu tiên
    - Image: resize giữ tỉ lệ
    - DOCX/PPTX/khác: trả None (FE dùng icon mặc định theo loại file)
    """
    path = Path(file_path)
    if not path.exists():
        return None

    out_path = THUMBNAIL_DIR / f"{doc_id}.jpg"
    group = get_file_group(mime_type)

    try:
        if group == "pdf":
            from pdf2image import convert_from_path
            images = convert_from_path(str(path), first_page=1, last_page=1, dpi=150)
            if images:
                images[0].convert("RGB").save(str(out_path), "JPEG", quality=85)
                return f"storage/thumbnails/{doc_id}.jpg"

        elif group == "image":
            from PIL import Image
            img = Image.open(path).convert("RGB")
            img.thumbnail((800, 800))
            img.save(str(out_path), "JPEG", quality=85)
            return f"storage/thumbnails/{doc_id}.jpg"

    except Exception as e:
        logger.error(f"create_thumbnail lỗi [{file_path}]: {e}")

    return None


def _ocr_pdf(path: Path) -> str:
    """OCR tối đa 5 trang đầu của PDF scan"""
    try:
        import pytesseract
        from pdf2image import convert_from_path
        images = convert_from_path(str(path), first_page=1, last_page=5, dpi=200)
        return "\n".join(
            pytesseract.image_to_string(img, lang="vie+eng") for img in images
        )
    except Exception as e:
        logger.error(f"OCR PDF lỗi: {e}")
        return ""


def _ocr_image(path: Path) -> str:
    try:
        import pytesseract
        from PIL import Image
        return pytesseract.image_to_string(Image.open(path), lang="vie+eng")
    except Exception as e:
        logger.error(f"OCR image lỗi: {e}")
        return ""